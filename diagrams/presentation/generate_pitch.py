#!/usr/bin/env python3
"""Generate Nepal360 Pitch Presentation (.pptx)"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand Colors ──
EMERALD       = RGBColor(0x05, 0x96, 0x69)  # #059669
EMERALD_DARK  = RGBColor(0x04, 0x78, 0x57)  # #047857
EMERALD_LIGHT = RGBColor(0x10, 0xB9, 0x81)  # #10b981
EMERALD_PALE  = RGBColor(0xD1, 0xFA, 0xE5)  # #d1fae5
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
DARK_SLATE    = RGBColor(0x0F, 0x17, 0x2A)  # #0f172a
LIGHT_GRAY    = RGBColor(0xF8, 0xFA, 0xFC)  # #f8fafc
MID_GRAY      = RGBColor(0x64, 0x74, 0x8B)  # #64748b
RED_ACCENT    = RGBColor(0xEF, 0x44, 0x44)  # #ef4444
AMBER         = RGBColor(0xF5, 0x9E, 0x0B)  # #f59e0b

REPORT_IMAGES = os.path.join(os.path.dirname(__file__), '..', '..', 'report_images')
OUTPUT_PATH   = os.path.join(os.path.dirname(__file__), 'nepal360_pitch.pptx')

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


# ── Helper Functions ──

def add_solid_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    s.line.fill.background()
    return s

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=DARK_SLATE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=DARK_SLATE, bullet_color=EMERALD, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = spacing
        p.level = 0
        # Bullet character
        p.text = f"\u2022  {item}"
    return txBox

def add_section_header(slide, title, subtitle=None):
    """Dark top bar with title"""
    add_shape(slide, 0, 0, W, Inches(1.4), DARK_SLATE)
    # Emerald accent line
    add_shape(slide, 0, Inches(1.4), W, Inches(0.06), EMERALD)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
                 title, font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.5),
                     subtitle, font_size=16, color=EMERALD_LIGHT)

def add_feature_card(slide, left, top, width, height, title, desc, icon_text=""):
    """Card with emerald top border"""
    # Card background
    card = add_shape(slide, left, top, width, height, LIGHT_GRAY)
    card.shadow.inherit = False
    # Top accent
    add_shape(slide, left, top, width, Inches(0.06), EMERALD)
    # Icon circle
    if icon_text:
        circle = add_shape(slide, left + Inches(0.3), top + Inches(0.25),
                          Inches(0.5), Inches(0.5), EMERALD, MSO_SHAPE.OVAL)
        circle.text_frame.paragraphs[0].text = icon_text
        circle.text_frame.paragraphs[0].font.size = Pt(16)
        circle.text_frame.paragraphs[0].font.color.rgb = WHITE
        circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Title
    add_text_box(slide, left + Inches(0.2), top + Inches(0.85), width - Inches(0.4), Inches(0.4),
                 title, font_size=14, color=DARK_SLATE, bold=True)
    # Description
    add_text_box(slide, left + Inches(0.2), top + Inches(1.2), width - Inches(0.4), Inches(0.6),
                 desc, font_size=11, color=MID_GRAY)

def add_image_safe(slide, img_name, left, top, width=None, height=None):
    """Add image from report_images if it exists"""
    path = os.path.join(REPORT_IMAGES, img_name)
    if os.path.exists(path):
        kwargs = {'left': left, 'top': top}
        if width: kwargs['width'] = width
        if height: kwargs['height'] = height
        slide.shapes.add_picture(path, **kwargs)
        return True
    else:
        # Placeholder rectangle
        w = width or Inches(5)
        h = height or Inches(3)
        ph = add_shape(slide, left, top, w, h, LIGHT_GRAY)
        ph.text_frame.paragraphs[0].text = f"[{img_name}]"
        ph.text_frame.paragraphs[0].font.size = Pt(14)
        ph.text_frame.paragraphs[0].font.color.rgb = MID_GRAY
        ph.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        return False

def add_placeholder_screenshot(slide, left, top, width, height, label):
    """Placeholder for UI screenshots"""
    border = add_shape(slide, left, top, width, height, LIGHT_GRAY)
    border.line.color.rgb = EMERALD
    border.line.width = Pt(2)
    # Center text
    add_text_box(slide, left, top + height / 2 - Inches(0.3), width, Inches(0.6),
                 f"[ {label} ]", font_size=18, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

def add_slide_number(slide, num, total):
    add_text_box(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
                 f"{num}/{total}", font_size=10, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)


TOTAL_SLIDES = 17

# ═══════════════════════════════════════════════════════════════
# SLIDE 1: Title Slide
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_solid_bg(slide, DARK_SLATE)

# Large emerald accent shape - angled
add_shape(slide, 0, 0, Inches(5.5), H, EMERALD_DARK)
# Lighter overlay triangle effect
pts = [Inches(4), 0, Inches(6.5), 0, Inches(4), H]
add_shape(slide, Inches(4), 0, Inches(2.5), H, EMERALD)

# Logo text
add_text_box(slide, Inches(6.5), Inches(1.5), Inches(6), Inches(1.2),
             "Nepal360", font_size=64, color=WHITE, bold=True)

# Tagline
add_text_box(slide, Inches(6.5), Inches(2.8), Inches(6), Inches(0.6),
             "Empowering Change Across Nepal", font_size=24, color=EMERALD_LIGHT)

# Divider line
add_shape(slide, Inches(6.5), Inches(3.6), Inches(3), Inches(0.04), EMERALD)

# Presenter info
add_text_box(slide, Inches(6.5), Inches(3.9), Inches(6), Inches(0.5),
             "Presented by Nischal", font_size=18, color=WHITE)
add_text_box(slide, Inches(6.5), Inches(4.4), Inches(6), Inches(0.5),
             "Full-Stack Developer", font_size=14, color=MID_GRAY)

# Nepal text on left panel
add_text_box(slide, Inches(0.8), Inches(2.5), Inches(4), Inches(1),
             "360\u00b0", font_size=120, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.8), Inches(5.2), Inches(4), Inches(0.5),
             "Crowdfunding Platform", font_size=16, color=EMERALD_PALE, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 1, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2: About Me
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)
add_section_header(slide, "About Me", "The developer behind Nepal360")

# Profile placeholder circle
circle = add_shape(slide, Inches(1.2), Inches(2.2), Inches(2.5), Inches(2.5),
                   EMERALD_PALE, MSO_SHAPE.OVAL)
circle.text_frame.paragraphs[0].text = "Photo"
circle.text_frame.paragraphs[0].font.size = Pt(20)
circle.text_frame.paragraphs[0].font.color.rgb = EMERALD
circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Info
add_text_box(slide, Inches(4.5), Inches(2.2), Inches(7), Inches(0.6),
             "Nischal", font_size=36, color=DARK_SLATE, bold=True)
add_text_box(slide, Inches(4.5), Inches(2.9), Inches(7), Inches(0.5),
             "Full-Stack Developer", font_size=20, color=EMERALD)

details = [
    "BCA Student at Itahari Namuna College",
    "Passionate about building impactful tech solutions for Nepal",
    "Experienced in React, Node.js, TypeScript, PostgreSQL",
    "Nepal360 is my final year capstone project",
]
add_bullet_list(slide, Inches(4.5), Inches(3.6), Inches(7.5), Inches(3),
                details, font_size=16, color=DARK_SLATE)

add_slide_number(slide, 2, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3: Problem Statement
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)
add_section_header(slide, "The Problem in Nepal", "Why Nepal needs a better crowdfunding platform")

problems = [
    ("No Organized Platform", "Nepal lacks a dedicated, trustworthy crowdfunding infrastructure for social causes"),
    ("Trust Deficit", "Donors are skeptical due to lack of transparency and accountability in charitable giving"),
    ("No KYC Verification", "Campaign creators are not verified, leading to potential fraud and misuse of funds"),
    ("Unregulated Sector", "Nepal's charity and social funding sector remains largely unregulated and fragmented"),
]

for i, (title, desc) in enumerate(problems):
    col = i % 2
    row = i // 2
    left = Inches(0.8) + col * Inches(6.2)
    top = Inches(1.9) + row * Inches(2.3)

    # Card
    card = add_shape(slide, left, top, Inches(5.6), Inches(2.0), LIGHT_GRAY)
    # Red accent on left
    add_shape(slide, left, top, Inches(0.08), Inches(2.0), RED_ACCENT)
    # Warning icon
    icon = add_shape(slide, left + Inches(0.3), top + Inches(0.3),
                     Inches(0.5), Inches(0.5), RED_ACCENT, MSO_SHAPE.OVAL)
    icon.text_frame.paragraphs[0].text = "!"
    icon.text_frame.paragraphs[0].font.size = Pt(20)
    icon.text_frame.paragraphs[0].font.color.rgb = WHITE
    icon.text_frame.paragraphs[0].font.bold = True
    icon.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(slide, left + Inches(1.0), top + Inches(0.25), Inches(4.2), Inches(0.4),
                 title, font_size=18, color=DARK_SLATE, bold=True)
    add_text_box(slide, left + Inches(1.0), top + Inches(0.75), Inches(4.2), Inches(1.0),
                 desc, font_size=14, color=MID_GRAY)

add_slide_number(slide, 3, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4: Impact of the Problem
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)
add_section_header(slide, "Impact of the Problem", "What the lack of infrastructure has led to")

impacts = [
    ("\u26a0", "Fraud & Misuse", "Donated funds are misused with no accountability, eroding public trust in charitable giving"),
    ("\u2716", "No Tracking", "Donors cannot track where their money goes or see the real impact of their contributions"),
    ("\u2709", "Hard to Reach Donors", "Beneficiaries and cause organizers struggle to connect with potential donors effectively"),
    ("\u2302", "Rural Communities Left Behind", "Remote and rural communities remain underserved due to lack of digital fundraising access"),
]

for i, (icon, title, desc) in enumerate(impacts):
    left = Inches(0.6) + i * Inches(3.1)
    top = Inches(2.0)

    # Tall card
    card = add_shape(slide, left, top, Inches(2.8), Inches(4.5), LIGHT_GRAY)
    add_shape(slide, left, top, Inches(2.8), Inches(0.06), EMERALD)

    # Icon
    add_text_box(slide, left, top + Inches(0.3), Inches(2.8), Inches(0.6),
                 icon, font_size=36, color=EMERALD, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), top + Inches(1.1), Inches(2.4), Inches(0.5),
                 title, font_size=17, color=DARK_SLATE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), top + Inches(1.7), Inches(2.4), Inches(2.5),
                 desc, font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 4, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5: The Solution
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)
add_section_header(slide, "Our Solution", "How Nepal360 addresses these challenges")

solutions = [
    ("\u2713", "KYC-Verified Campaigns",
     "Every campaign creator must complete identity verification before launching a campaign, ensuring authenticity"),
    ("\u2261", "Full Transparency & Tracking",
     "Real-time donation tracking, progress bars, donor lists, and detailed financial breakdowns for every campaign"),
    ("\u2699", "AI-Powered Intelligence",
     "Machine learning predictions for fundraising success, smart campaign recommendations for donors"),
    ("\u2316", "Map-Based Discovery",
     "Interactive district-level map of Nepal for discovering campaigns by location, making local causes visible"),
]

for i, (icon, title, desc) in enumerate(solutions):
    top = Inches(1.8) + i * Inches(1.3)
    # Emerald circle icon
    c = add_shape(slide, Inches(0.8), top + Inches(0.1), Inches(0.6), Inches(0.6),
                  EMERALD, MSO_SHAPE.OVAL)
    c.text_frame.paragraphs[0].text = icon
    c.text_frame.paragraphs[0].font.size = Pt(16)
    c.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Title + desc
    add_text_box(slide, Inches(1.8), top, Inches(10), Inches(0.4),
                 title, font_size=20, color=DARK_SLATE, bold=True)
    add_text_box(slide, Inches(1.8), top + Inches(0.45), Inches(10), Inches(0.7),
                 desc, font_size=14, color=MID_GRAY)

add_slide_number(slide, 5, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6: Introducing Nepal360 - Features
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)
add_section_header(slide, "Introducing Nepal360", "A comprehensive crowdfunding platform for Nepal")

features = [
    ("\u270e", "Campaign\nManagement", "Create, manage\n& track campaigns"),
    ("\u2611", "KYC\nVerification", "Identity verified\ncampaign creators"),
    ("\u2b50", "Khalti\nPayments", "Secure digital\npayment gateway"),
    ("\u2606", "Leaderboards\n& Badges", "Gamified donor\nengagement"),
    ("\u2699", "AI\nPredictions", "ML-powered funding\nforecasts"),
    ("\u266a", "Real-Time\nNotifications", "Instant updates\non campaigns"),
    ("\u2709", "Certificate\nGeneration", "Auto-generated\ndonor certificates"),
    ("\u2316", "Interactive\nMap", "District-level\ncampaign discovery"),
]

for i, (icon, title, desc) in enumerate(features):
    col = i % 4
    row = i // 4
    left = Inches(0.6) + col * Inches(3.1)
    top = Inches(1.9) + row * Inches(2.6)

    add_feature_card(slide, left, top, Inches(2.8), Inches(2.2), title, desc, icon)

add_slide_number(slide, 6, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDES 7-10: UI Screenshots (placeholders)
# ═══════════════════════════════════════════════════════════════
screenshot_slides = [
    ("UI Preview — Home Page", "Homepage with featured campaigns, stats, and hero section"),
    ("UI Preview — Campaign Discovery", "Browse campaigns with filters, search, and map view"),
    ("UI Preview — Campaign Detail + Donation", "Campaign details, donation form, progress tracking, and donor list"),
    ("UI Preview — Admin Dashboard", "Admin panel with KYC verification, campaign management, and analytics"),
]

for idx, (title, subtitle) in enumerate(screenshot_slides):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_solid_bg(slide, WHITE)
    add_section_header(slide, title, subtitle)

    # Large placeholder for screenshot
    add_placeholder_screenshot(slide, Inches(1.0), Inches(1.8), Inches(11.3), Inches(5.2),
                               f"Insert {title.split('—')[1].strip()} Screenshot Here")

    add_slide_number(slide, 7 + idx, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDE 11: Competitor Analysis - Global
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)
add_section_header(slide, "Competitor Analysis — Global", "How Nepal360 compares to global platforms")

# Competitor images + info cards
competitors = [
    ("gofundme.png", "GoFundMe", "Global leader in personal fundraising.\nNo Nepal focus, no KYC, high fees."),
    ("globalgiving.png", "GlobalGiving", "Non-profit focused platform.\nLimited to vetted organizations only."),
    ("ketto.png", "Ketto", "India-based crowdfunding.\nNo Nepal presence, no local payments."),
]

for i, (img, name, desc) in enumerate(competitors):
    left = Inches(0.5) + i * Inches(4.2)
    top = Inches(1.8)

    # Card
    card = add_shape(slide, left, top, Inches(3.8), Inches(5.0), LIGHT_GRAY)
    add_shape(slide, left, top, Inches(3.8), Inches(0.06), EMERALD)

    # Competitor screenshot
    add_image_safe(slide, img, left + Inches(0.2), top + Inches(0.3),
                   width=Inches(3.4), height=Inches(2.2))

    # Name
    add_text_box(slide, left + Inches(0.2), top + Inches(2.7), Inches(3.4), Inches(0.4),
                 name, font_size=20, color=DARK_SLATE, bold=True)
    # Description
    add_text_box(slide, left + Inches(0.2), top + Inches(3.2), Inches(3.4), Inches(1.5),
                 desc, font_size=13, color=MID_GRAY)

add_slide_number(slide, 11, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDE 12: Competitor Analysis - Nepal
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)
add_section_header(slide, "Competitor Analysis — Nepal", "The local competitive landscape")

# SahayogiHaat card
add_image_safe(slide, "sahayogihaat.png", Inches(0.8), Inches(1.9),
               width=Inches(4.5), height=Inches(3.0))

add_text_box(slide, Inches(0.8), Inches(5.1), Inches(4.5), Inches(0.5),
             "SahayogiHaat — Only Notable Local Competitor", font_size=16, color=DARK_SLATE, bold=True)

# What they lack - comparison
lacks = [
    "No KYC verification system",
    "No AI-powered predictions",
    "No gamification (badges, leaderboards)",
    "No interactive map discovery",
    "No item/in-kind donations",
    "No recurring donation support",
    "No certificate generation",
    "No real-time notifications",
]

add_text_box(slide, Inches(6.0), Inches(1.9), Inches(6.5), Inches(0.5),
             "What they lack (Nepal360 advantages):", font_size=18, color=EMERALD, bold=True)

for i, item in enumerate(lacks):
    top = Inches(2.6) + i * Inches(0.5)
    # Red X
    add_text_box(slide, Inches(6.0), top, Inches(0.4), Inches(0.4),
                 "\u2717", font_size=16, color=RED_ACCENT, bold=True)
    add_text_box(slide, Inches(6.5), top, Inches(6), Inches(0.4),
                 item, font_size=15, color=DARK_SLATE)

add_slide_number(slide, 12, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDE 13: Market Stats & Opportunity
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)
add_section_header(slide, "Market Opportunity", "Why now is the right time for Nepal360")

stats = [
    ("~27%", "of GDP from Remittances", "Nepal has one of the highest\nremittance-to-GDP ratios globally"),
    ("56%+", "Internet Penetration", "16.5M internet users with\nrapidly growing connectivity"),
    ("86%", "Wallet-to-Population Ratio", "Khalti & eSewa driving\nmassive digital payment adoption"),
    ("12th", "Climate Risk Index", "Highly vulnerable to floods,\nlandslides & earthquakes"),
]

for i, (stat, label, desc) in enumerate(stats):
    left = Inches(0.5) + i * Inches(3.2)
    top = Inches(1.9)

    # Stat card
    card = add_shape(slide, left, top, Inches(2.9), Inches(3.5), LIGHT_GRAY)
    add_shape(slide, left, top, Inches(2.9), Inches(0.06), EMERALD)

    # Big number
    add_text_box(slide, left, top + Inches(0.3), Inches(2.9), Inches(0.8),
                 stat, font_size=42, color=EMERALD, bold=True, alignment=PP_ALIGN.CENTER)
    # Label
    add_text_box(slide, left + Inches(0.2), top + Inches(1.2), Inches(2.5), Inches(0.5),
                 label, font_size=16, color=DARK_SLATE, bold=True, alignment=PP_ALIGN.CENTER)
    # Description
    add_text_box(slide, left + Inches(0.2), top + Inches(1.8), Inches(2.5), Inches(1.2),
                 desc, font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Key insight bar
add_shape(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.2), EMERALD_PALE)
add_text_box(slide, Inches(0.8), Inches(5.95), Inches(12), Inches(0.8),
             "\u2736  Nepal's unique combination of high remittance inflow, growing digital payments, and disaster vulnerability creates a massive opportunity for a transparent, trusted crowdfunding platform.",
             font_size=14, color=EMERALD_DARK)

add_slide_number(slide, 13, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDE 14: System Architecture
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)
add_section_header(slide, "System Architecture", "Technical overview of the Nepal360 platform")

add_image_safe(slide, "architecture.png", Inches(1.5), Inches(1.8),
               width=Inches(10.3), height=Inches(5.2))

add_slide_number(slide, 14, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDE 15: Revenue Model
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)
add_section_header(slide, "Revenue Model", "How Nepal360 generates sustainable income")

revenue_streams = [
    ("\u00a4", "Platform Fee", "2-5% fee on successful donations",
     "Primary revenue stream. Small, transparent fee applied only when campaigns reach their goals."),
    ("\u2605", "Premium Features", "Enhanced campaign tools for creators",
     "Priority placement, advanced analytics, custom branding, and extended campaign durations."),
    ("\u2302", "Corporate Partnerships", "CSR & corporate matching programs",
     "Partner with businesses for matched donations, corporate-sponsored campaigns, and CSR integration."),
    ("\u2606", "Sponsored Campaigns", "Featured and promoted campaigns",
     "Organizations can sponsor campaign visibility, badges, and featured placement on the platform."),
]

for i, (icon, title, subtitle, desc) in enumerate(revenue_streams):
    col = i % 2
    row = i // 2
    left = Inches(0.6) + col * Inches(6.2)
    top = Inches(1.9) + row * Inches(2.5)

    # Card
    card = add_shape(slide, left, top, Inches(5.8), Inches(2.2), LIGHT_GRAY)
    add_shape(slide, left, top, Inches(0.08), Inches(2.2), EMERALD)

    # Icon
    add_text_box(slide, left + Inches(0.3), top + Inches(0.3), Inches(0.6), Inches(0.6),
                 icon, font_size=28, color=EMERALD, alignment=PP_ALIGN.CENTER)
    # Title
    add_text_box(slide, left + Inches(1.0), top + Inches(0.25), Inches(4.5), Inches(0.4),
                 title, font_size=20, color=DARK_SLATE, bold=True)
    # Subtitle
    add_text_box(slide, left + Inches(1.0), top + Inches(0.65), Inches(4.5), Inches(0.3),
                 subtitle, font_size=14, color=EMERALD)
    # Description
    add_text_box(slide, left + Inches(1.0), top + Inches(1.1), Inches(4.5), Inches(0.9),
                 desc, font_size=13, color=MID_GRAY)

add_slide_number(slide, 15, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDE 16: Project Demo
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_SLATE)

# Large emerald shape
add_shape(slide, 0, Inches(2.0), W, Inches(3.5), EMERALD_DARK)
add_shape(slide, 0, Inches(2.0), W, Inches(0.06), EMERALD_LIGHT)
add_shape(slide, 0, Inches(5.44), W, Inches(0.06), EMERALD_LIGHT)

# Play icon circle
circle = add_shape(slide, Inches(5.9), Inches(2.5), Inches(1.5), Inches(1.5),
                   EMERALD, MSO_SHAPE.OVAL)
circle.text_frame.paragraphs[0].text = "\u25b6"
circle.text_frame.paragraphs[0].font.size = Pt(48)
circle.text_frame.paragraphs[0].font.color.rgb = WHITE
circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

add_text_box(slide, 0, Inches(4.2), W, Inches(0.8),
             "Live Demo", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 0, Inches(5.8), W, Inches(0.5),
             "Nepal360 Platform Walkthrough", font_size=18, color=EMERALD_LIGHT, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 16, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════
# SLIDE 17: Thank You
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_SLATE)

# Emerald accent bars
add_shape(slide, 0, Inches(0), Inches(0.15), H, EMERALD)
add_shape(slide, Inches(13.183), Inches(0), Inches(0.15), H, EMERALD)

add_text_box(slide, 0, Inches(1.5), W, Inches(1.2),
             "Thank You!", font_size=64, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(5.5), Inches(2.8), Inches(2.3), Inches(0.04), EMERALD)

add_text_box(slide, 0, Inches(3.2), W, Inches(0.6),
             "Nepal360 — Empowering Change Across Nepal", font_size=22, color=EMERALD_LIGHT,
             alignment=PP_ALIGN.CENTER)

# Contact info
add_text_box(slide, 0, Inches(4.2), W, Inches(0.5),
             "Nischal", font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 0, Inches(4.7), W, Inches(0.5),
             "BCA Student | Itahari Namuna College", font_size=16, color=MID_GRAY,
             alignment=PP_ALIGN.CENTER)

# Bottom branding bar
add_shape(slide, 0, Inches(6.5), W, Inches(1.0), EMERALD_DARK)
add_text_box(slide, 0, Inches(6.65), W, Inches(0.5),
             "Questions?", font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 17, TOTAL_SLIDES)


# ── Save ──
prs.save(OUTPUT_PATH)
print(f"Presentation saved to: {OUTPUT_PATH}")
print(f"Total slides: {TOTAL_SLIDES}")
